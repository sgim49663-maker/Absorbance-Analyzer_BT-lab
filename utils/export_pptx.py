"""
분석 결과를 PowerPoint 파일로 내보내기.
- Slide 1: 타이틀
- Slide 2: 바 차트 이미지
- Slide 3: 통계 테이블
"""
import io
import tempfile
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from utils.visualization import create_bar_chart


def generate_pptx(result: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 타이틀 ──
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 배경 색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xD9)

    # 타이틀 텍스트
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Absorbance Analysis Results"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "96-well Plate Absorbance Assay"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p2.alignment = PP_ALIGN.CENTER

    # ── Slide 2: 바 차트 ──
    slide2 = prs.slides.add_slide(slide_layout)
    fig = create_bar_chart(result, figsize=(12, 6))

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        fig.savefig(tmp.name, dpi=200, bbox_inches="tight", facecolor="white")
        tmp_path = tmp.name

    slide2.shapes.add_picture(
        tmp_path,
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(11.5), height=Inches(6.2),
    )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def _style_cell(cell, bold=False, font_size=Pt(11),
                bg_color=None, font_color=RGBColor(0x33, 0x33, 0x33)):
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = font_color

    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if bg_color:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg_color.__str__()})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
